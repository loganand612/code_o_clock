#!/usr/bin/env python3
"""
PPT Generator Service for Course Summary Presentations

This module provides functionality to generate PowerPoint presentations
that summarize course content with effective visual presentations.
"""

from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import base64
from datetime import datetime


class PPTGenerator:
    """Generates PowerPoint presentations for course summaries."""
    
    def __init__(self):
        """Initialize the PPT generator with default styling."""
        # Define color scheme
        self.primary_color = RGBColor(25, 118, 210)  # Blue #1976d2
        self.secondary_color = RGBColor(76, 175, 80)  # Green #4caf50
        self.accent_color = RGBColor(255, 152, 0)   # Orange #ff9800
        self.text_color = RGBColor(33, 33, 33)       # Dark grey #212121
        self.light_grey = RGBColor(245, 245, 245)   # Light grey #f5f5f5
        
        # Font sizes
        self.title_font_size = Pt(44)
        self.heading_font_size = Pt(32)
        self.subheading_font_size = Pt(24)
        self.body_font_size = Pt(18)
        self.small_font_size = Pt(14)
    
    def create_title_slide(self, prs: Presentation, course_title: str, course_summary: str = ""):
        """Create the title slide for the presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = course_title
        title_frame = title.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = self.title_font_size
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self.primary_color
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"Course Summary\nGenerated on {datetime.now().strftime('%B %d, %Y')}"
        subtitle_frame = subtitle.text_frame
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.font.size = self.subheading_font_size
        subtitle_paragraph.font.color.rgb = self.text_color
        
        if course_summary:
            # Add a second paragraph with course summary
            p = subtitle_frame.add_paragraph()
            p.text = course_summary
            p.font.size = self.body_font_size
            p.font.color.rgb = self.text_color
            p.space_after = Pt(12)
    
    def create_overview_slide(self, prs: Presentation, modules: List[Dict[str, Any]]):
        """Create a course overview slide showing all modules."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide title
        title = slide.shapes.title
        title.text = "Course Overview"
        title_frame = title.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = self.heading_font_size
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self.primary_color
        
        # Add content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.margin_left = Inches(0.5)
        content_frame.margin_right = Inches(0.5)
        
        # Add module list
        for i, module in enumerate(modules, 1):
            if i > 1:  # Add spacing between modules
                p = content_frame.add_paragraph()
                p.text = ""
                p.space_after = Pt(6)
            
            p = content_frame.add_paragraph()
            p.text = f"Module {i}: {module['title']}"
            p.font.size = self.body_font_size
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            p.space_after = Pt(4)
            
            # Add lesson count and estimated time
            lesson_count = len(module.get('lessons', []))
            estimated_time = lesson_count * 30  # 30 minutes per lesson
            
            p = content_frame.add_paragraph()
            p.text = f"   • {lesson_count} lessons • ~{estimated_time} minutes"
            p.font.size = self.small_font_size
            p.font.color.rgb = self.text_color
            p.space_after = Pt(2)
    
    def create_module_summary_slide(self, prs: Presentation, module: Dict[str, Any], module_num: int):
        """Create a summary slide for a specific module."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide title
        title = slide.shapes.title
        title.text = f"Module {module_num}: {module['title']}"
        title_frame = title.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = self.heading_font_size
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self.primary_color
        
        # Add content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.margin_left = Inches(0.5)
        content_frame.margin_right = Inches(0.5)
        
        # Add module description if available
        if module.get('description'):
            p = content_frame.add_paragraph()
            p.text = module['description']
            p.font.size = self.body_font_size
            p.font.color.rgb = self.text_color
            p.space_after = Pt(12)
            p.space_before = Pt(6)
        
        # Add lessons
        p = content_frame.add_paragraph()
        p.text = "Key Lessons:"
        p.font.size = self.subheading_font_size
        p.font.bold = True
        p.font.color.rgb = self.secondary_color
        p.space_after = Pt(8)
        
        for i, lesson in enumerate(module.get('lessons', []), 1):
            p = content_frame.add_paragraph()
            p.text = f"{i}. {lesson['title']}"
            p.font.size = self.body_font_size
            p.font.bold = True
            p.font.color.rgb = self.text_color
            p.space_after = Pt(4)
            
            # Add lesson summary if available
            if lesson.get('summary'):
                p = content_frame.add_paragraph()
                p.text = f"   {lesson['summary']}"
                p.font.size = self.small_font_size
                p.font.color.rgb = self.text_color
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(8)
                p.level = 1  # Indent this paragraph
    
    def create_key_concepts_slide(self, prs: Presentation, modules: List[Dict[str, Any]]):
        """Create a slide highlighting key concepts from all modules."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide title
        title = slide.shapes.title
        title.text = "Key Concepts & Takeaways"
        title_frame = title.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = self.heading_font_size
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self.primary_color
        
        # Add content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.margin_left = Inches(0.5)
        content_frame.margin_right = Inches(0.5)
        
        # Extract key concepts from lessons
        key_concepts = []
        for module in modules:
            for lesson in module.get('lessons', []):
                if lesson.get('summary'):
                    # Extract key terms from summary (simple approach)
                    summary = lesson['summary']
                    # Add lesson title as key concept
                    key_concepts.append({
                        'concept': lesson['title'],
                        'module': module['title']
                    })
        
        # Add key concepts to slide
        if key_concepts:
            p = content_frame.add_paragraph()
            p.text = "Core Learning Objectives:"
            p.font.size = self.subheading_font_size
            p.font.bold = True
            p.font.color.rgb = self.secondary_color
            p.space_after = Pt(8)
            
            for i, concept in enumerate(key_concepts[:8], 1):  # Limit to 8 concepts
                p = content_frame.add_paragraph()
                p.text = f"• {concept['concept']}"
                p.font.size = self.body_font_size
                p.font.color.rgb = self.text_color
                p.space_after = Pt(4)
        
        # Add completion message
        p = content_frame.add_paragraph()
        p.text = ""
        p.space_after = Pt(12)
        
        p = content_frame.add_paragraph()
        p.text = "✓ Course Completed Successfully!"
        p.font.size = self.body_font_size
        p.font.bold = True
        p.font.color.rgb = self.secondary_color
        p.space_after = Pt(6)
    
    def generate_course_ppt(self, course_data: Dict[str, Any]) -> bytes:
        """
        Generate a complete PowerPoint presentation for the course.
        
        Args:
            course_data: Dictionary containing course information with structure:
                {
                    "course": "Course Title",
                    "modules": [
                        {
                            "title": "Module Title",
                            "description": "Module description",
                            "lessons": [
                                {
                                    "title": "Lesson Title",
                                    "summary": "Lesson summary"
                                }
                            ]
                        }
                    ]
                }
        
        Returns:
            bytes: The generated PowerPoint file as bytes
        """
        # Create presentation
        prs = Presentation()
        
        # Get course information
        course_title = course_data.get('course', 'Course Summary')
        modules = course_data.get('modules', [])
        
        # Create title slide
        self.create_title_slide(prs, course_title, f"Complete course with {len(modules)} modules")
        
        # Create overview slide
        if modules:
            self.create_overview_slide(prs, modules)
        
        # Create module summary slides
        for i, module in enumerate(modules, 1):
            self.create_module_summary_slide(prs, module, i)
        
        # Create key concepts slide
        if modules:
            self.create_key_concepts_slide(prs, modules)
        
        # Save to bytes
        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        
        return ppt_bytes.getvalue()
    
    def generate_ppt_base64(self, course_data: Dict[str, Any]) -> str:
        """
        Generate a PowerPoint presentation and return as base64 string.
        
        Args:
            course_data: Dictionary containing course information
        
        Returns:
            str: Base64 encoded PowerPoint file
        """
        ppt_bytes = self.generate_course_ppt(course_data)
        return base64.b64encode(ppt_bytes).decode('utf-8')


# Global PPT generator instance
ppt_generator = PPTGenerator()