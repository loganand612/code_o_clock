"""
PowerPoint Generator for Course Content
Generates structured PowerPoint presentations from course data
"""

import os
import uuid
import json
import requests
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

class PowerPointGenerator:
    """Generates PowerPoint presentations from course content"""
    
    def __init__(self, gemini_api_key: str = None):
        self.gemini_api_key = gemini_api_key or os.environ.get('GOOGLE_API_KEY')
        self.gemini_model = "gemini-2.0-flash-exp"
        self.api_url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.gemini_model}:generateContent?key={self.gemini_api_key}"
        
        # Create downloads directory if it doesn't exist
        self.download_folder = 'downloads'
        if not os.path.exists(self.download_folder):
            os.makedirs(self.download_folder)
    
    def generate_from_course(self, course_data: Dict[str, Any]) -> str:
        """
        Generate PowerPoint from course data
        Returns the filename of the generated presentation
        """
        if not course_data.get('generatedCourse'):
            raise ValueError("No course data available for PowerPoint generation")
        
        course = course_data['generatedCourse']
        course_title = course.get('course', 'Generated Course')
        modules = course.get('modules', [])
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        self._add_title_slide(prs, course_title, course_data.get('courseDetails', {}))
        
        # Add module slides
        for module_idx, module in enumerate(modules, 1):
            self._add_module_slide(prs, module, module_idx)
            
            # Add lesson slides
            for lesson_idx, lesson in enumerate(module.get('lessons', []), 1):
                self._add_lesson_slide(prs, lesson, module_idx, lesson_idx)
        
        # Add summary slide
        self._add_summary_slide(prs, course_title, len(modules))
        
        # Save presentation
        filename = f"course-{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(self.download_folder, filename)
        prs.save(filepath)
        
        return filename
    
    def generate_from_text(self, text: str, title: str = "Generated Presentation") -> str:
        """
        Generate PowerPoint from raw text using AI structuring
        """
        # Get structured slide data from AI
        slide_data = self._get_slide_structure_from_ai(text)
        
        if not slide_data:
            raise ValueError("Failed to generate structured slides from text")
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        self._add_custom_title_slide(prs, title)
        
        # Add content slides
        for slide_info in slide_data:
            self._add_content_slide(prs, slide_info)
        
        # Save presentation
        filename = f"presentation-{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(self.download_folder, filename)
        prs.save(filepath)
        
        return filename
    
    def _add_title_slide(self, prs: Presentation, course_title: str, course_details: Dict):
        """Add a professional title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = course_title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Add subtitle with course details
        subtitle = slide.placeholders[1]
        subtitle_text = subtitle.text_frame
        subtitle_text.clear()
        
        if course_details.get('creatorName'):
            p = subtitle_text.paragraphs[0]
            p.text = f"Created by: {course_details['creatorName']}"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(102, 102, 102)
        
        if course_details.get('difficulty'):
            p = subtitle_text.add_paragraph()
            p.text = f"Level: {course_details['difficulty']}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(102, 102, 102)
    
    def _add_custom_title_slide(self, prs: Presentation, title: str):
        """Add a simple title slide"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    def _add_module_slide(self, prs: Presentation, module: Dict, module_num: int):
        """Add a module overview slide"""
        slide_layout = prs.slide_layouts[1]  # Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Module title
        title_shape = slide.shapes.title
        title_shape.text = f"Module {module_num}: {module.get('title', 'Untitled Module')}"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Module content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add lesson overview
        p = text_frame.paragraphs[0]
        p.text = "Lessons in this module:"
        p.font.size = Pt(24)
        p.font.bold = True
        
        for i, lesson in enumerate(module.get('lessons', []), 1):
            p = text_frame.add_paragraph()
            p.text = f"• {lesson.get('title', 'Untitled Lesson')}"
            p.font.size = Pt(20)
            p.level = 1
    
    def _add_lesson_slide(self, prs: Presentation, lesson: Dict, module_num: int, lesson_num: int):
        """Add a lesson content slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Lesson title
        title_shape = slide.shapes.title
        title_shape.text = f"Module {module_num}.{lesson_num}: {lesson.get('title', 'Untitled Lesson')}"
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Lesson content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add summary
        if lesson.get('summary'):
            p = text_frame.paragraphs[0]
            p.text = "Summary:"
            p.font.size = Pt(20)
            p.font.bold = True
            
            p = text_frame.add_paragraph()
            p.text = lesson['summary']
            p.font.size = Pt(18)
            p.level = 1
        
        # Add key points from detail
        if lesson.get('detail'):
            p = text_frame.add_paragraph()
            p.text = "Key Points:"
            p.font.size = Pt(20)
            p.font.bold = True
            
            # Extract key points from detail (simplified)
            detail_text = lesson['detail']
            sentences = detail_text.split('.')[:5]  # Take first 5 sentences
            
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"• {sentence.strip()}"
                    p.font.size = Pt(16)
                    p.level = 1
    
    def _add_content_slide(self, prs: Presentation, slide_info: Dict):
        """Add a content slide from AI-generated structure"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = slide_info.get('title', 'Untitled Slide')
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for bullet in slide_info.get('bullets', []):
            p = text_frame.paragraphs[0] if not text_frame.paragraphs else text_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.level = 0
    
    def _add_summary_slide(self, prs: Presentation, course_title: str, module_count: int):
        """Add a summary slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Course Summary"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Summary content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = f"Congratulations! You have completed:"
        p.font.size = Pt(24)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"• {course_title}"
        p.font.size = Pt(20)
        p.level = 1
        
        p = text_frame.add_paragraph()
        p.text = f"• {module_count} comprehensive modules"
        p.font.size = Pt(20)
        p.level = 1
        
        p = text_frame.add_paragraph()
        p.text = "Thank you for learning with us!"
        p.font.size = Pt(18)
        p.font.italic = True
    
    def _get_slide_structure_from_ai(self, text: str) -> List[Dict]:
        """Get structured slide data from AI (using Gemini)"""
        if not self.gemini_api_key:
            print("Warning: No Gemini API key available, using fallback structure")
            return self._fallback_slide_structure(text)
        
        schema = {
            "type": "ARRAY",
            "description": "An array of slides, where each slide is an object containing a title and a list of bullet points.",
            "items": {
                "type": "OBJECT",
                "properties": {
                    "title": {"type": "STRING", "description": "The concise title for the PowerPoint slide."},
                    "bullets": {
                        "type": "ARRAY",
                        "items": {"type": "STRING", "description": "A single, concise bullet point."},
                        "description": "A list of key points for the slide body."
                    }
                },
                "required": ["title", "bullets"]
            }
        }
        
        system_prompt = (
            "You are an expert presentation designer. Analyze the provided text and organize it into "
            "logical PowerPoint slides. Each slide should have a compelling title and 3-5 concise, "
            "informative bullet points. Return the output as a JSON array matching the specified schema."
        )
        
        payload = {
            "contents": [{"parts": [{"text": text}]}],
            "systemInstruction": {"parts": [{"text": system_prompt}]},
            "generationConfig": {
                "responseMimeType": "application/json",
                "responseSchema": schema,
            }
        }
        
        try:
            response = requests.post(
                self.api_url, 
                headers={"Content-Type": "application/json"}, 
                data=json.dumps(payload), 
                timeout=30
            )
            response.raise_for_status()
            
            result = response.json()
            
            if (result.get('candidates') and 
                result['candidates'][0].get('content') and 
                result['candidates'][0]['content'].get('parts')):
                
                json_text = result['candidates'][0]['content']['parts'][0]['text']
                slide_data = json.loads(json_text)
                print(f"AI generated {len(slide_data)} slides successfully.")
                return slide_data
                
        except Exception as e:
            print(f"AI generation failed: {e}")
        
        return self._fallback_slide_structure(text)
    
    def _fallback_slide_structure(self, text: str) -> List[Dict]:
        """Fallback slide structure when AI is not available"""
        # Simple text splitting as fallback
        paragraphs = text.split('\n\n')
        slides = []
        
        for i, paragraph in enumerate(paragraphs[:10]):  # Max 10 slides
            if paragraph.strip():
                sentences = paragraph.split('.')[:5]  # Max 5 bullets per slide
                bullets = [s.strip() for s in sentences if s.strip()]
                
                slides.append({
                    "title": f"Slide {i + 1}",
                    "bullets": bullets
                })
        
        return slides

# Global instance
pptx_generator = PowerPointGenerator()
