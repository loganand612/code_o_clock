import fitz  # PyMuPDF
import docx
from pptx import Presentation
from pytube import YouTube
import re
import requests
from bs4 import BeautifulSoup

def extract_text_from_pdf(file_stream):
    text = ""
    with fitz.open(stream=file_stream.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(file_stream):
    doc = docx.Document(file_stream)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_pptx(file_stream):
    prs = Presentation(file_stream)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_youtube(url):
    try:
        # Create YouTube object
        yt = YouTube(url)
        
        # Get video title and description
        title = yt.title
        description = yt.description
        
        # Get captions
        caption_text = ""
        try:
            captions = yt.captions
            if 'en' in captions:
                caption_text = captions['en'].generate_srt_captions()
            elif len(captions) > 0:
                # Get first available caption if English is not available
                first_caption = list(captions.values())[0]
                caption_text = first_caption.generate_srt_captions()
        except Exception as e:
            print(f"Could not get captions: {str(e)}")
        
        # Clean caption text (remove timecodes and numbers)
        if caption_text:
            # Remove SRT formatting
            caption_text = re.sub(r'\d+\n\d{2}:\d{2}:\d{2},\d{3} --> \d{2}:\d{2}:\d{2},\d{3}\n', '', caption_text)
            caption_text = caption_text.replace('\n\n', ' ').replace('\n', ' ')
        
        # Combine all text
        all_text = f"Title: {title}\n\nDescription: {description}"
        if caption_text:
            all_text += f"\n\nTranscript: {caption_text}"
            
        return all_text.strip()
            
    except Exception as e:
        print(f"Error extracting YouTube content: {str(e)}")
        return ""

def extract_text_from_website(url):
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()  # Raise an exception for bad status codes
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove unwanted elements
        for element in soup(['script', 'style', 'header', 'footer', 'nav', 'ads', 'iframe']):
            element.decompose()
            
        # Try to find the main content
        main_content = None
        for tag in ['article', 'main', '[role="main"]', '.content', '#content', '.post-content']:
            main_content = soup.select_one(tag)
            if main_content:
                break
                
        if main_content:
            text = main_content.get_text()
        else:
            text = soup.get_text()
            
        # Clean up the text
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = ' '.join(chunk for chunk in chunks if chunk)
        
        # Remove multiple spaces and clean up
        text = ' '.join(text.split())
        return text.strip()
        
    except Exception as e:
        print(f"Error extracting website text: {str(e)}")
        return ""